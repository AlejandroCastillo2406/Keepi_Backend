from typing import Dict, Any, List, Optional
import tempfile
import os
from PIL import Image
import easyocr
from datetime import datetime, timedelta
import re
from app.services.bedrock_service import BedrockService

class DocumentAnalysisService:
    """Servicio para análisis automático de documentos usando AI"""
    
    def __init__(self):
        # Inicializar EasyOCR para español e inglés
        self.easyocr_reader = easyocr.Reader(['es', 'en'])
        # Inicializar Bedrock para análisis con Claude 3 Haiku
        self.bedrock_service = BedrockService()
    
    async def analyze_document(self, content: bytes, content_type: str, filename: str) -> Dict[str, Any]:
        """Analizar documento y extraer información automáticamente"""
        try:
            # Extraer texto del documento
            extracted_text = await self._extract_text(content, content_type, filename)
            
            # Clasificar documento
            suggested_category = await self._classify_document(extracted_text, filename)
            
            # Si requiere clasificación manual, devolver respuesta especial
            if suggested_category == "MANUAL_CLASSIFICATION_REQUIRED":
                return {
                    "suggested_category": "MANUAL_CLASSIFICATION_REQUIRED",
                    "confidence_score": 0.0,
                    "extracted_text": "",
                    "metadata": {},
                    "tags": ["manual_classification_required"],
                    "expiry_date": None,
                    "document_number": None,
                    "organization": None,
                    "processing_time_ms": 0,
                    "ai_model_version": "1.0.0",
                    "manual_classification_message": "No pudimos clasificarlo de manera adecuada, ¿a qué categoría corresponde?"
                }
            
            # Extraer metadatos
            metadata = await self._extract_metadata(extracted_text)
            
            # Generar tags
            tags = await self._generate_tags(extracted_text, suggested_category)
            
            # Calcular confianza
            confidence_score = await self._calculate_confidence(extracted_text, suggested_category)
            
            # Extraer fecha de vencimiento
            expiry_date = await self._extract_expiry_date(extracted_text)
            
            # Extraer número de documento
            document_number = await self._extract_document_number(extracted_text)
            
            # Extraer organización
            organization = await self._extract_organization(extracted_text)
            
            return {
                "suggested_category": suggested_category,
                "confidence_score": confidence_score,
                "extracted_text": extracted_text,
                "metadata": metadata,
                "tags": tags,
                "expiry_date": expiry_date,
                "document_number": document_number,
                "organization": organization,
                "processing_time_ms": 0,  # TODO: Implementar medición de tiempo
                "ai_model_version": "1.0.0"
            }
            
        except Exception as e:
            print(f"Error analizando documento: {e}")
            # Retornar análisis básico en caso de error
            return {
                "suggested_category": "Documento",
                "confidence_score": 0.1,
                "extracted_text": "",
                "metadata": {},
                "tags": ["error"],
                "expiry_date": None,
                "document_number": None,
                "organization": None,
                "processing_time_ms": 0,
                "ai_model_version": "1.0.0"
            }
    
    async def _extract_text(self, content: bytes, content_type: str, filename: str) -> str:
        """Extraer texto del documento según su tipo"""
        try:
            if content_type.startswith('image/'):
                # Procesar imagen con estrategia: Textract -> EasyOCR -> fallback manual
                return await self._extract_text_from_image(content, filename)
            elif content_type == 'application/pdf':
                # Usar estrategia: Textract asíncrono -> EasyOCR -> fallback manual
                return await self._extract_text_from_pdf(content, filename)
            elif content_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                # Procesar Word: python-docx -> openpyxl -> fallback manual
                return await self._extract_text_from_word(content, filename)
            elif content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
                # Procesar Excel: openpyxl -> fallback manual
                return await self._extract_text_from_excel(content, filename)
            elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                # Procesar PowerPoint: python-pptx -> fallback manual
                return await self._extract_text_from_powerpoint(content, filename)
            else:
                # Para otros tipos, intentar decodificar como texto
                try:
                    return content.decode('utf-8')
                except:
                    return f"Archivo: {filename}"
                    
        except Exception as e:
            print(f"Error extrayendo texto: {e}")
            return f"Error extrayendo texto: {filename}"
    
    async def _extract_text_from_image(self, content: bytes, filename: str) -> str:
        """Extraer texto de imagen usando estrategia: Textract -> EasyOCR -> fallback manual"""
        try:
            # Estrategia 1: AWS Textract (principal)
            try:
                from app.services.aws_service import AWSService
                aws_service = AWSService()
                result = await aws_service._extract_text_from_image(content)
                extracted_text = result.get('text', '')
                if extracted_text and len(extracted_text.strip()) > 50:
                    print(f"✅ Textract extrajo texto de imagen {filename}: {len(extracted_text)} caracteres")
                    return extracted_text
                else:
                    print(f"⚠️ Textract no pudo extraer texto suficiente de {filename}")
            except Exception as e:
                print(f"❌ Error en Textract para imagen {filename}: {e}")
            
            # Estrategia 2: EasyOCR (fallback)
            try:
                print(f"🔄 Intentando EasyOCR para imagen {filename}...")
            # Crear archivo temporal
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name
            
            try:
                    # Usar EasyOCR
                    results = self.easyocr_reader.readtext(temp_file_path)
                    text = ' '.join([result[1] for result in results])
                    
                    if text and len(text.strip()) > 20:
                        print(f"✅ EasyOCR extrajo texto de imagen {filename}: {len(text)} caracteres")
                        return text.strip()
                    else:
                        print(f"⚠️ EasyOCR no pudo extraer texto suficiente de {filename}")
                        
                finally:
                    # Limpiar archivo temporal
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                
            except Exception as e:
                print(f"❌ Error en EasyOCR para imagen {filename}: {e}")
            
            # Estrategia 3: Fallback manual
            print(f"⚠️ No se pudo extraer texto de imagen {filename}, requiere clasificación manual")
            return "MANUAL_CLASSIFICATION_REQUIRED"
                
        except Exception as e:
            print(f"Error procesando imagen {filename}: {e}")
            return "MANUAL_CLASSIFICATION_REQUIRED"
    
    async def _extract_text_from_pdf(self, content: bytes, filename: str) -> str:
        """Extraer texto de PDF usando estrategia: Textract asíncrono -> EasyOCR -> fallback manual"""
        try:
            # Estrategia 1: AWS Textract asíncrono (principal)
            try:
                from app.services.aws_service import AWSService
                aws_service = AWSService()
                result = await aws_service.extract_text_from_document(content, filename, 'application/pdf')
                extracted_text = result.get('text', '')
                if extracted_text and len(extracted_text.strip()) > 100:
                    print(f"✅ Textract asíncrono extrajo texto de PDF {filename}: {len(extracted_text)} caracteres")
                    return extracted_text
                else:
                    print(f"⚠️ Textract asíncrono no pudo extraer texto suficiente de {filename}")
            except Exception as e:
                print(f"❌ Error en Textract asíncrono para PDF {filename}: {e}")
            
            # Estrategia 2: Clasificación manual (fallback)
            print(f"⚠️ PDF {filename} requiere clasificación manual")
            return "MANUAL_CLASSIFICATION_REQUIRED"
                
        except Exception as e:
            print(f"Error procesando PDF {filename}: {e}")
            return "MANUAL_CLASSIFICATION_REQUIRED"
    
    async def _extract_text_from_word(self, content: bytes, filename: str) -> str:
        """Extraer texto de Word usando estrategia: python-docx -> fallback manual"""
        try:
            # Estrategia 1: python-docx (principal)
            try:
                from docx import Document
                import io
                
                doc = Document(io.BytesIO(content))
                all_text = []
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        all_text.append(paragraph.text)
                
                # Extraer texto de tablas
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                all_text.append(cell.text)
                
                text = '\n'.join(all_text)
                
                if text and len(text.strip()) > 50:
                    print(f"✅ python-docx extrajo texto de Word {filename}: {len(text)} caracteres")
                    return text.strip()
                else:
                    print(f"⚠️ python-docx no pudo extraer texto suficiente de {filename}")
                    
            except Exception as e:
                print(f"❌ Error en python-docx para Word {filename}: {e}")
            
            # Estrategia 2: Fallback manual
            print(f"⚠️ No se pudo extraer texto de Word {filename}, requiere clasificación manual")
            return "MANUAL_CLASSIFICATION_REQUIRED"
                
        except Exception as e:
            print(f"Error procesando Word {filename}: {e}")
            return "MANUAL_CLASSIFICATION_REQUIRED"
    
    async def _extract_text_from_excel(self, content: bytes, filename: str) -> str:
        """Extraer texto de Excel usando estrategia: openpyxl -> fallback manual"""
        try:
            # Estrategia 1: openpyxl (principal)
            try:
                from openpyxl import load_workbook
                import io
                
                workbook = load_workbook(io.BytesIO(content))
                all_text = []
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows():
                        row_text = []
                        for cell in row:
                            if cell.value and str(cell.value).strip():
                                row_text.append(str(cell.value))
                        if row_text:
                            all_text.append(' '.join(row_text))
                
                text = '\n'.join(all_text)
                
                if text and len(text.strip()) > 50:
                    print(f"✅ openpyxl extrajo texto de Excel {filename}: {len(text)} caracteres")
                    return text.strip()
                else:
                    print(f"⚠️ openpyxl no pudo extraer texto suficiente de {filename}")
                    
            except Exception as e:
                print(f"❌ Error en openpyxl para Excel {filename}: {e}")
            
            # Estrategia 2: Fallback manual
            print(f"⚠️ No se pudo extraer texto de Excel {filename}, requiere clasificación manual")
            return "MANUAL_CLASSIFICATION_REQUIRED"
                
        except Exception as e:
            print(f"Error procesando Excel {filename}: {e}")
            return "MANUAL_CLASSIFICATION_REQUIRED"
    
    async def _extract_text_from_powerpoint(self, content: bytes, filename: str) -> str:
        """Extraer texto de PowerPoint usando estrategia: python-pptx -> fallback manual"""
        try:
            # Estrategia 1: python-pptx (principal)
            try:
                from pptx import Presentation
                import io
                
                prs = Presentation(io.BytesIO(content))
                all_text = []
                
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            all_text.append(shape.text)
                
                text = '\n'.join(all_text)
                
                if text and len(text.strip()) > 50:
                    print(f"✅ python-pptx extrajo texto de PowerPoint {filename}: {len(text)} caracteres")
                return text.strip()
                else:
                    print(f"⚠️ python-pptx no pudo extraer texto suficiente de {filename}")
                
            except Exception as e:
                print(f"❌ Error en python-pptx para PowerPoint {filename}: {e}")
            
            # Estrategia 2: Fallback manual
            print(f"⚠️ No se pudo extraer texto de PowerPoint {filename}, requiere clasificación manual")
            return "MANUAL_CLASSIFICATION_REQUIRED"
                
        except Exception as e:
            print(f"Error procesando PowerPoint {filename}: {e}")
            return "MANUAL_CLASSIFICATION_REQUIRED"
    
    async def _classify_document(self, text: str, filename: str) -> str:
        """Clasificar documento usando Amazon Bedrock con Claude 3 Haiku"""
        try:
            # Si requiere clasificación manual, devolver señal especial
            if text == "MANUAL_CLASSIFICATION_REQUIRED":
                return "MANUAL_CLASSIFICATION_REQUIRED"
            
            # Usar Bedrock con Claude 3 Haiku para clasificación
            bedrock_result = await self.bedrock_service.analyze_document_content(text, filename)
            category = bedrock_result.get('category', 'Documento')
            
            return category
            
        except Exception as e:
            print(f"Error en clasificación con Bedrock: {e}")
            # Fallback: clasificación básica por extensión si Bedrock falla
            if filename.endswith('.pdf'):
                return "Documento PDF"
            elif filename.endswith(('.jpg', '.jpeg', '.png')):
                return "Imagen"
            else:
                return "Documento"
    
    async def _extract_metadata(self, text: str) -> Dict[str, Any]:
        """Extraer metadatos del texto"""
        metadata = {}
        
        # Buscar fechas
        date_patterns = [
            r'\d{1,2}/\d{1,2}/\d{4}',  # DD/MM/YYYY
            r'\d{1,2}-\d{1,2}-\d{4}',  # DD-MM-YYYY
            r'\d{4}-\d{1,2}-\d{1,2}',  # YYYY-MM-DD
        ]
        
        dates_found = []
        for pattern in date_patterns:
            dates = re.findall(pattern, text)
            dates_found.extend(dates)
        
        if dates_found:
            metadata['fechas_encontradas'] = list(set(dates_found))
        
        # Buscar montos
        amount_patterns = [
            r'S/\.\s*\d+[,.]?\d*',  # S/. 150.00
            r'\$\s*\d+[,.]?\d*',    # $ 150.00
            r'\d+[,.]?\d*\s*USD',   # 150.00 USD
        ]
        
        amounts_found = []
        for pattern in amount_patterns:
            amounts = re.findall(pattern, text)
            amounts_found.extend(amounts)
        
        if amounts_found:
            metadata['montos'] = list(set(amounts_found))
        
        # Buscar números de documento
        doc_patterns = [
            r'[A-Z]{1,3}-\d{4,}',  # F001-2024
            r'[A-Z]{1,3}\d{4,}',   # F0012024
            r'\d{4,}',              # 2024001
        ]
        
        doc_numbers = []
        for pattern in doc_patterns:
            numbers = re.findall(pattern, text)
            doc_numbers.extend(numbers)
        
        if doc_numbers:
            metadata['numeros_documento'] = list(set(doc_numbers))
        
        return metadata
    
    async def _generate_tags(self, text: str, category: str) -> List[str]:
        """Generar tags automáticamente usando Bedrock"""
        try:
            tags = [category.lower()]
            
            # Usar Bedrock para generar tags dinámicos
            bedrock_result = await self.bedrock_service.analyze_document_content(text, "temp")
            
            # Agregar tags básicos basados en el texto
            text_lower = text.lower()
            basic_tags = {
                "urgente": ["urgente", "inmediato", "asap"],
                "importante": ["importante", "crítico", "prioritario"],
                "confidencial": ["confidencial", "privado", "secreto"],
                "borrador": ["borrador", "draft", "temporal"]
            }
            
            for tag, keywords in basic_tags.items():
                if any(keyword in text_lower for keyword in keywords):
                    tags.append(tag)
            
            # Agregar tags basados en palabras clave del texto
            text_lower = text.lower()
            keyword_tags = []
            
            # Palabras clave comunes para generar tags
            if any(word in text_lower for word in ["certificado", "diploma", "título", "académico"]):
                keyword_tags.append("académico")
            if any(word in text_lower for word in ["contrato", "nómina", "trabajo", "empleo"]):
                keyword_tags.append("laboral")
            if any(word in text_lower for word in ["receta", "análisis", "médico", "salud"]):
                keyword_tags.append("médico")
            if any(word in text_lower for word in ["factura", "estado", "cuenta", "pago"]):
                keyword_tags.append("financiero")
            if any(word in text_lower for word in ["dni", "pasaporte", "licencia", "identidad"]):
                keyword_tags.append("identificación")
            if any(word in text_lower for word in ["matrícula", "seguro", "itv", "auto"]):
                keyword_tags.append("vehículo")
            if any(word in text_lower for word in ["alquiler", "escritura", "casa", "hogar"]):
                keyword_tags.append("vivienda")
            
            tags.extend(keyword_tags)
            
            return list(set(tags))[:10]  # Máximo 10 tags únicos
            
        except Exception as e:
            print(f"Error generando tags con Bedrock: {e}")
            # Fallback: tags básicos
            return [category.lower(), "documento"]
    
    async def _calculate_confidence(self, text: str, category: str) -> float:
        """Calcular nivel de confianza de la clasificación usando Bedrock"""
        try:
            if not text or text.strip() == "":
                return 0.1
            
            # Usar Bedrock para obtener confianza real
            bedrock_result = await self.bedrock_service.analyze_document_content(text, "temp")
            confidence = bedrock_result.get('confidence', 0.5)
            
            # Normalizar confianza entre 0 y 1
            return round(min(max(confidence, 0.0), 1.0), 2)
            
        except Exception as e:
            print(f"Error calculando confianza con Bedrock: {e}")
            # Fallback: confianza básica basada en longitud del texto
            text_length = len(text)
            length_confidence = min(text_length / 1000, 1.0)
            return round(length_confidence, 2)
    
    async def _extract_expiry_date(self, text: str) -> Optional[str]:
        """Extraer fecha de vencimiento del texto usando Bedrock"""
        try:
            # Usar Bedrock para extraer fecha de vencimiento
            bedrock_result = await self.bedrock_service.analyze_document_content(text, "temp")
            expiry_date = bedrock_result.get('expiry_date')
            
            if expiry_date and expiry_date != "null":
                return expiry_date
            
            # Fallback: buscar patrones de fecha de vencimiento
        expiry_patterns = [
            r'venc[ei]miento[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})',
            r'expir[ae][:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})',
            r'validez[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})',
            r'vigencia[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})'
        ]
        
        for pattern in expiry_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1)
        
        return None
            
        except Exception as e:
            print(f"Error extrayendo fecha de vencimiento con Bedrock: {e}")
        return None
    
    async def _extract_document_number(self, text: str) -> Optional[str]:
        """Extraer número de documento del texto"""
        # Buscar patrones de número de documento
        doc_patterns = [
            r'[Nn]úmero[:\s]*([A-Z]{1,3}[-]?\d{4,})',
            r'[Cc]ódigo[:\s]*([A-Z]{1,3}[-]?\d{4,})',
            r'[Rr]eferencia[:\s]*([A-Z]{1,3}[-]?\d{4,})',
            r'[Ii]D[:\s]*([A-Z]{1,3}[-]?\d{4,})'
        ]
        
        for pattern in doc_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1)
        
        return None
    
    async def _extract_organization(self, text: str) -> Optional[str]:
        """Extraer nombre de organización del texto"""
        # Buscar patrones de organización
        org_patterns = [
            r'[Ee]mpresa[:\s]*([A-Z][a-z\s]+)',
            r'[Cc]ompañía[:\s]*([A-Z][a-z\s]+)',
            r'[Ii]nstitución[:\s]*([A-Z][a-z\s]+)',
            r'[Oo]rganización[:\s]*([A-Z][a-z\s]+)'
        ]
        
        for pattern in org_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1).strip()
        
        return None
