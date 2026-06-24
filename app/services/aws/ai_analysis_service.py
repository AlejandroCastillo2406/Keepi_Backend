import logging
import re
from typing import Any, Dict, List, Optional

from sqlalchemy.orm import Session

from app.core.config import settings
from app.services.aws.bedrock_service import BedrockService
from app.services.subscription import SubscriptionService

logger = logging.getLogger(__name__)

_UNLIMITED_ANALYSIS_CHECK = {
    "can_analyze": True,
    "plan": "unlimited",
    "analysis_used": 0,
    "analysis_remaining": 999999,
    "needs_subscription": False,
}


class DocumentAnalysisService:

    def __init__(self):

        self.bedrock_service = BedrockService()

        self.subscription_service = SubscriptionService()

    async def _analysis_quota(self, user_id: str, db: Session) -> Dict[str, Any]:
        if not settings.require_subscription_for_document_analysis:
            return dict(_UNLIMITED_ANALYSIS_CHECK)
        return await self.subscription_service.check_analysis_limit(user_id, db)

    async def _track_analysis_usage(self, user_id: str, db: Session) -> None:
        if settings.require_subscription_for_document_analysis:
            await self.subscription_service.increment_analysis_usage(user_id, db)

    async def analyze_document(
        self,
        content: bytes,
        content_type: str,
        filename: str,
        user_id: str,
        db: Session,
        existing_category_names: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        try:

            analysis_check = await self._analysis_quota(user_id, db)

            if (
                settings.require_subscription_for_document_analysis
                and not analysis_check["can_analyze"]
            ):
                return {
                    "suggested_category": "SUBSCRIPTION_REQUIRED",
                    "confidence_score": 0.0,
                    "extracted_text": "",
                    "metadata": {},
                    "tags": ["subscription_required"],
                    "expiry_date": None,
                    "processing_time_ms": 0,
                    "ai_model_version": "1.0.0",
                    "subscription_required_message": "Has alcanzado el límite de 2 análisis gratuitos. Suscríbete para obtener análisis ilimitados.",
                    "subscription_info": {
                        "current_plan": analysis_check["plan"],
                        "analysis_used": analysis_check["analysis_used"],
                        "analysis_remaining": analysis_check["analysis_remaining"],
                        "needs_subscription": analysis_check["needs_subscription"],
                    },
                }

            extracted_text = await self._extract_text(content, content_type, filename)

            if content_type.startswith("image/") and (
                extracted_text == "MANUAL_CLASSIFICATION_REQUIRED"
                or not extracted_text
                or len(extracted_text.strip()) < 30
            ):
                bedrock_vision = await self.bedrock_service.analyze_image_for_category(
                    content,
                    filename,
                    content_type,
                    existing_folder_names=existing_category_names or [],
                )
                suggested = (bedrock_vision.get("category") or "").strip()
                confidence = float(bedrock_vision.get("confidence") or 0)
                if (
                    suggested
                    and suggested != "MANUAL_CLASSIFICATION_REQUIRED"
                    and confidence >= 0.25
                ):
                    await self._track_analysis_usage(user_id, db)
                    tags = list(bedrock_vision.get("tags") or [])
                    if suggested.lower() not in [t.lower() for t in tags]:
                        tags.insert(0, suggested.lower())
                    return {
                        "suggested_category": suggested,
                        "confidence_score": confidence,
                        "extracted_text": "",
                        "metadata": {},
                        "tags": tags[:10],
                        "expiry_date": bedrock_vision.get("expiry_date"),
                        "recommended_name": bedrock_vision.get("recommended_name"),
                        "processing_time_ms": 0,
                        "ai_model_version": "1.0.0",
                        "subscription_info": {
                            "current_plan": analysis_check["plan"],
                            "analysis_used": analysis_check["analysis_used"] + 1,
                            "analysis_remaining": max(
                                0, analysis_check["analysis_remaining"] - 1
                            ),
                            "needs_subscription": False,
                        },
                    }

                return {
                    "suggested_category": "MANUAL_CLASSIFICATION_REQUIRED",
                    "confidence_score": 0.0,
                    "extracted_text": "",
                    "metadata": {},
                    "tags": ["manual_classification_required"],
                    "expiry_date": None,
                    "processing_time_ms": 0,
                    "ai_model_version": "1.0.0",
                    "manual_classification_message": "No pudimos clasificarlo de manera adecuada, ¿a qué categoría corresponde?",
                }

            if (
                extracted_text == "MANUAL_CLASSIFICATION_REQUIRED"
                or not extracted_text
                or len(extracted_text.strip()) < 30
            ):
                return {
                    "suggested_category": "MANUAL_CLASSIFICATION_REQUIRED",
                    "confidence_score": 0.0,
                    "extracted_text": "",
                    "metadata": {},
                    "tags": ["manual_classification_required"],
                    "expiry_date": None,
                    "processing_time_ms": 0,
                    "ai_model_version": "1.0.0",
                    "manual_classification_message": "No pudimos clasificarlo de manera adecuada, ¿a qué categoría corresponde?",
                }

            bedrock_result = await self.bedrock_service.analyze_document_content(
                extracted_text,
                filename,
                existing_folder_names=existing_category_names or [],
            )
            suggested_category = bedrock_result.get("category", "Documento")
            if suggested_category == "MANUAL_CLASSIFICATION_REQUIRED":
                return {
                    "suggested_category": "MANUAL_CLASSIFICATION_REQUIRED",
                    "confidence_score": 0.0,
                    "extracted_text": "",
                    "metadata": {},
                    "tags": ["manual_classification_required"],
                    "expiry_date": None,
                    "processing_time_ms": 0,
                    "ai_model_version": "1.0.0",
                    "manual_classification_message": "No pudimos clasificarlo de manera adecuada, ¿a qué categoría corresponde?",
                }

            confidence_score = bedrock_result.get("confidence", 0.5)
            expiry_date = bedrock_result.get(
                "expiry_date"
            ) or self._extract_expiry_date_regex_only(extracted_text)
            recommended_name = bedrock_result.get("recommended_name")
            tags = list(bedrock_result.get("tags") or [])
            if suggested_category.lower() not in [t.lower() for t in tags]:
                tags.insert(0, suggested_category.lower())
            tags = self._merge_keyword_tags(extracted_text, tags)

            metadata = await self._extract_metadata(extracted_text)

            await self._track_analysis_usage(user_id, db)

            return {
                "suggested_category": suggested_category,
                "confidence_score": confidence_score,
                "extracted_text": extracted_text,
                "metadata": metadata,
                "tags": tags[:10],
                "expiry_date": expiry_date,
                "recommended_name": recommended_name,
                "processing_time_ms": 0,
                "ai_model_version": "1.0.0",
                "subscription_info": {
                    "current_plan": analysis_check["plan"],
                    "analysis_used": analysis_check["analysis_used"] + 1,
                    "analysis_remaining": max(
                        0, analysis_check["analysis_remaining"] - 1
                    ),
                    "needs_subscription": False,
                },
            }

        except Exception as e:
            logger.warning("Error analizando documento: %s", e)

            return {
                "suggested_category": "Documento",
                "confidence_score": 0.1,
                "extracted_text": "",
                "metadata": {},
                "tags": ["error"],
                "expiry_date": None,
                "processing_time_ms": 0,
                "ai_model_version": "1.0.0",
            }

    async def _extract_text(
        self, content: bytes, content_type: str, filename: str
    ) -> str:
        try:
            if content_type.startswith("image/"):

                return await self._extract_text_from_image(content, filename)
            elif content_type == "application/pdf":

                return await self._extract_text_from_pdf(content, filename)
            elif content_type in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            ]:

                return await self._extract_text_from_word(content, filename)
            elif content_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ]:

                return await self._extract_text_from_excel(content, filename)
            elif (
                content_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):

                return await self._extract_text_from_powerpoint(content, filename)
            else:

                try:
                    return content.decode("utf-8")
                except:
                    return f"Archivo: {filename}"

        except Exception as e:
            logger.warning("Error extrayendo texto: %s", e)
            return f"Error extrayendo texto: {filename}"

    async def _extract_text_from_image(self, content: bytes, filename: str) -> str:
        try:

            try:
                from app.services.aws.aws_service import AWSService

                aws_service = AWSService()
                result = await aws_service._extract_text_from_image(content)
                extracted_text = result.get("text", "")
                if extracted_text and len(extracted_text.strip()) > 50:
                    logger.info(
                        "Textract extrajo texto de imagen %s: %d caracteres",
                        filename,
                        len(extracted_text),
                    )
                    return extracted_text
                else:
                    logger.warning(
                        "Textract no extrajo texto suficiente de %s", filename
                    )
            except Exception as e:
                logger.warning("Error en Textract para imagen %s: %s", filename, e)

            logger.warning("Imagen %s requiere clasificación manual", filename)
            return "MANUAL_CLASSIFICATION_REQUIRED"

        except Exception as e:
            logger.warning("Error procesando imagen %s: %s", filename, e)
            return "MANUAL_CLASSIFICATION_REQUIRED"

    async def _extract_text_from_pdf(self, content: bytes, filename: str) -> str:
        try:

            try:
                from app.services.aws.aws_service import AWSService

                aws_service = AWSService()
                result = await aws_service.extract_text_from_document(
                    content, filename, "application/pdf"
                )
                extracted_text = result.get("text", "")
                if extracted_text and len(extracted_text.strip()) > 100:
                    logger.info(
                        "Textract extrajo texto de PDF %s: %d caracteres",
                        filename,
                        len(extracted_text),
                    )
                    return extracted_text
                else:
                    logger.warning(
                        "Textract no extrajo texto suficiente de PDF %s", filename
                    )
            except Exception as e:
                logger.warning("Error en Textract para PDF %s: %s", filename, e)

            logger.warning("PDF %s requiere clasificación manual", filename)
            return "MANUAL_CLASSIFICATION_REQUIRED"

        except Exception as e:
            logger.warning("Error procesando PDF %s: %s", filename, e)
            return "MANUAL_CLASSIFICATION_REQUIRED"

    async def _extract_text_from_word(self, content: bytes, filename: str) -> str:
        try:

            try:
                import io

                from docx import Document

                doc = Document(io.BytesIO(content))
                all_text = []

                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        all_text.append(paragraph.text)

                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                all_text.append(cell.text)

                text = "\n".join(all_text)

                if text and len(text.strip()) > 50:
                    logger.info(
                        "python-docx extrajo texto de Word %s: %d caracteres",
                        filename,
                        len(text),
                    )
                    return text.strip()
                else:
                    logger.warning(
                        "python-docx no extrajo texto suficiente de %s", filename
                    )

            except Exception as e:
                logger.warning("Error en python-docx para Word %s: %s", filename, e)

            logger.warning("Word %s requiere clasificación manual", filename)
            return "MANUAL_CLASSIFICATION_REQUIRED"

        except Exception as e:
            logger.warning("Error procesando Word %s: %s", filename, e)
            return "MANUAL_CLASSIFICATION_REQUIRED"

    async def _extract_text_from_excel(self, content: bytes, filename: str) -> str:
        try:

            try:
                import io

                from openpyxl import load_workbook

                workbook = load_workbook(io.BytesIO(content))
                all_text = []

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows():
                        row_text = []
                        for cell in row:
                            if cell.value and str(cell.value).strip():
                                row_text.append(str(cell.value))
                        if row_text:
                            all_text.append(" ".join(row_text))

                text = "\n".join(all_text)

                if text and len(text.strip()) > 50:
                    logger.info(
                        "openpyxl extrajo texto de Excel %s: %d caracteres",
                        filename,
                        len(text),
                    )
                    return text.strip()
                else:
                    logger.warning(
                        "openpyxl no extrajo texto suficiente de %s", filename
                    )

            except Exception as e:
                logger.warning("Error en openpyxl para Excel %s: %s", filename, e)

            logger.warning("Excel %s requiere clasificación manual", filename)
            return "MANUAL_CLASSIFICATION_REQUIRED"

        except Exception as e:
            logger.warning("Error procesando Excel %s: %s", filename, e)
            return "MANUAL_CLASSIFICATION_REQUIRED"

    async def _extract_text_from_powerpoint(self, content: bytes, filename: str) -> str:
        try:

            try:
                import io

                from pptx import Presentation

                prs = Presentation(io.BytesIO(content))
                all_text = []

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            all_text.append(shape.text)

                text = "\n".join(all_text)

                if text and len(text.strip()) > 50:
                    logger.info(
                        "python-pptx extrajo texto de PowerPoint %s: %d caracteres",
                        filename,
                        len(text),
                    )
                    return text.strip()
                else:
                    logger.warning(
                        "python-pptx no extrajo texto suficiente de %s", filename
                    )

            except Exception as e:
                logger.warning(
                    "Error en python-pptx para PowerPoint %s: %s", filename, e
                )

            logger.warning("PowerPoint %s requiere clasificación manual", filename)
            return "MANUAL_CLASSIFICATION_REQUIRED"

        except Exception as e:
            logger.warning("Error procesando PowerPoint %s: %s", filename, e)
            return "MANUAL_CLASSIFICATION_REQUIRED"

    async def _classify_document(self, text: str, filename: str) -> str:
        try:

            if text == "MANUAL_CLASSIFICATION_REQUIRED":
                return "MANUAL_CLASSIFICATION_REQUIRED"

            bedrock_result = await self.bedrock_service.analyze_document_content(
                text, filename
            )
            category = bedrock_result.get("category", "Documento")

            return category

        except Exception as e:
            logger.warning("Error en clasificación con Bedrock: %s", e)

            if filename.endswith(".pdf"):
                return "Documento PDF"
            elif filename.endswith((".jpg", ".jpeg", ".png")):
                return "Imagen"
            else:
                return "Documento"

    async def _extract_metadata(self, text: str) -> Dict[str, Any]:
        metadata = {}

        date_patterns = [
            r"\d{1,2}/\d{1,2}/\d{4}",
            r"\d{1,2}-\d{1,2}-\d{4}",
            r"\d{4}-\d{1,2}-\d{1,2}",
        ]

        dates_found = []
        for pattern in date_patterns:
            dates = re.findall(pattern, text)
            dates_found.extend(dates)

        if dates_found:
            metadata["fechas_encontradas"] = list(set(dates_found))

        amount_patterns = [
            r"S/\.\s*\d+[,.]?\d*",
            r"\$\s*\d+[,.]?\d*",
            r"\d+[,.]?\d*\s*USD",
        ]

        amounts_found = []
        for pattern in amount_patterns:
            amounts = re.findall(pattern, text)
            amounts_found.extend(amounts)

        if amounts_found:
            metadata["montos"] = list(set(amounts_found))

        doc_patterns = [
            r"[A-Z]{1,3}-\d{4,}",
            r"[A-Z]{1,3}\d{4,}",
            r"\d{4,}",
        ]

        doc_numbers = []
        for pattern in doc_patterns:
            numbers = re.findall(pattern, text)
            doc_numbers.extend(numbers)

        if doc_numbers:
            metadata["numeros_documento"] = list(set(doc_numbers))

        return metadata

    async def _generate_tags(self, text: str, category: str) -> List[str]:
        try:
            tags = [category.lower()]

            bedrock_result = await self.bedrock_service.analyze_document_content(
                text, "temp"
            )

            text_lower = text.lower()
            basic_tags = {
                "urgente": ["urgente", "inmediato", "asap"],
                "importante": ["importante", "crítico", "prioritario"],
                "confidencial": ["confidencial", "privado", "secreto"],
                "borrador": ["borrador", "draft", "temporal"],
            }

            for tag, keywords in basic_tags.items():
                if any(keyword in text_lower for keyword in keywords):
                    tags.append(tag)

            text_lower = text.lower()
            keyword_tags = []

            if any(
                word in text_lower
                for word in ["certificado", "diploma", "título", "académico"]
            ):
                keyword_tags.append("académico")
            if any(
                word in text_lower
                for word in ["contrato", "nómina", "trabajo", "empleo"]
            ):
                keyword_tags.append("laboral")
            if any(
                word in text_lower for word in ["receta", "análisis", "médico", "salud"]
            ):
                keyword_tags.append("médico")
            if any(
                word in text_lower for word in ["factura", "estado", "cuenta", "pago"]
            ):
                keyword_tags.append("financiero")
            if any(
                word in text_lower
                for word in ["dni", "pasaporte", "licencia", "identidad"]
            ):
                keyword_tags.append("identificación")
            if any(
                word in text_lower for word in ["matrícula", "seguro", "itv", "auto"]
            ):
                keyword_tags.append("vehículo")
            if any(
                word in text_lower
                for word in ["alquiler", "escritura", "casa", "hogar"]
            ):
                keyword_tags.append("vivienda")

            tags.extend(keyword_tags)

            return list(set(tags))[:10]

        except Exception as e:
            logger.warning("Error generando tags con Bedrock: %s", e)

            return [category.lower(), "documento"]

    async def _calculate_confidence(self, text: str, category: str) -> float:
        try:
            if not text or text.strip() == "":
                return 0.1

            bedrock_result = await self.bedrock_service.analyze_document_content(
                text, "temp"
            )
            confidence = bedrock_result.get("confidence", 0.5)

            return round(min(max(confidence, 0.0), 1.0), 2)

        except Exception as e:
            logger.warning("Error calculando confianza con Bedrock: %s", e)

            text_length = len(text)
            length_confidence = min(text_length / 1000, 1.0)
            return round(length_confidence, 2)

    def _extract_expiry_date_regex_only(self, text: str) -> Optional[str]:
        expiry_patterns = [
            r"venc[ei]miento[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})",
            r"expir[ae][:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})",
            r"validez[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})",
            r"vigencia[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{4})",
        ]
        for pattern in expiry_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1)
        return None

    def _merge_keyword_tags(self, text: str, tags: List[str]) -> List[str]:
        text_lower = text.lower()
        basic_tags = {
            "urgente": ["urgente", "inmediato", "asap"],
            "importante": ["importante", "crítico", "prioritario"],
            "confidencial": ["confidencial", "privado", "secreto"],
            "borrador": ["borrador", "draft", "temporal"],
        }
        for tag, keywords in basic_tags.items():
            if tag not in tags and any(kw in text_lower for kw in keywords):
                tags.append(tag)
        keyword_map = [
            (["certificado", "diploma", "título", "académico"], "académico"),
            (["contrato", "nómina", "trabajo", "empleo"], "laboral"),
            (["receta", "análisis", "médico", "salud"], "médico"),
            (["factura", "estado", "cuenta", "pago"], "financiero"),
            (["dni", "pasaporte", "licencia", "identidad"], "identificación"),
            (["matrícula", "seguro", "itv", "auto"], "vehículo"),
            (["alquiler", "escritura", "casa", "hogar"], "vivienda"),
        ]
        for words, tag in keyword_map:
            if tag not in tags and any(w in text_lower for w in words):
                tags.append(tag)
        return list(dict.fromkeys(tags))[:10]

    async def _extract_expiry_date(self, text: str) -> Optional[str]:
        try:
            bedrock_result = await self.bedrock_service.analyze_document_content(
                text, "temp"
            )
            expiry_date = bedrock_result.get("expiry_date")
            if expiry_date and expiry_date != "null":
                return expiry_date
            return self._extract_expiry_date_regex_only(text)
        except Exception as e:
            logger.warning("Error extrayendo fecha de vencimiento con Bedrock: %s", e)
            return self._extract_expiry_date_regex_only(text)

    async def _extract_document_number(self, text: str) -> Optional[str]:

        doc_patterns = [
            r"[Nn]úmero[:\s]*([A-Z]{1,3}[-]?\d{4,})",
            r"[Cc]ódigo[:\s]*([A-Z]{1,3}[-]?\d{4,})",
            r"[Rr]eferencia[:\s]*([A-Z]{1,3}[-]?\d{4,})",
            r"[Ii]D[:\s]*([A-Z]{1,3}[-]?\d{4,})",
        ]

        for pattern in doc_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1)

        return None

    async def _extract_organization(self, text: str) -> Optional[str]:

        org_patterns = [
            r"[Ee]mpresa[:\s]*([A-Z][a-z\s]+)",
            r"[Cc]ompañía[:\s]*([A-Z][a-z\s]+)",
            r"[Ii]nstitución[:\s]*([A-Z][a-z\s]+)",
            r"[Oo]rganización[:\s]*([A-Z][a-z\s]+)",
        ]

        for pattern in org_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1).strip()

        return None
